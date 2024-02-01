from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Crear una presentación
prs = Presentation()

# Añadir diapositiva con título y contenido
slide_layout = prs.slide_layouts[5]  # Diseño con título y contenido
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Concepto y objetivo de la comunicación"

# Dar estilo al título
title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Color azul oscuro

# Añadir cuadro de texto con preguntas
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Dar estilo al cuadro de texto
tf.margin_left = Inches(0.5)
tf.margin_right = Inches(0.5)
tf.margin_top = Inches(0.5)
tf.margin_bottom = Inches(0.5)

# Dar estilo a los párrafos
for i, pregunta in enumerate(["¿Por qué me dirijo al público?", "¿Qué deseo conseguir?", "¿Qué deseo que las personas receptoras hagan o sientan después?"]):
    p = tf.add_paragraph()
    p.text = pregunta
    p.font.size = Pt(24)
    if i == 0:
        p.font.bold = True
    p.space_after = Pt(10)

# ... Continuar con el resto de las diapositivas y contenido ...

# Guardar la presentación
ruta_guardado = 'D:/mi_presentacion_bonita.pptx'
prs.save(ruta_guardado)
