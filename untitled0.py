# -*- coding: utf-8 -*-
"""
Created on Tue Aug 29 11:22:22 2023

@author: PC
"""


# Importar la biblioteca python-pptx
from pptx import Presentation
from pptx.util import Inches

# Crear una presentación vacía
prs = Presentation()

# Añadir una diapositiva con el diseño de título y contenido
slide_layout = prs.slide_layouts[1]

# Primera diapositiva: concepto y objetivo de la comunicación
slide = prs.slides.add_slide(slide_layout)

# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "Concepto y objetivo de la comunicación"

# Añadir un cuadro de texto con las preguntas
left = Inches(0.5) # Posición horizontal del cuadro de texto
top = Inches(1.5) # Posición vertical del cuadro de texto
width = Inches(9) # Ancho del cuadro de texto
height = Inches(5) # Alto del cuadro de texto
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Añadir las preguntas como párrafos al cuadro de texto
tf.add_paragraph("¿Por qué me dirijo al público?")
tf.add_paragraph("¿Qué deseo conseguir?")
tf.add_paragraph("¿Qué deseo que las personas receptoras hagan o sientan después?")

# Segunda diapositiva: la audiencia
slide = prs.slides.add_slide(slide_layout)

# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "La audiencia"

# Añadir un cuadro de texto con las preguntas
left = Inches(0.5) # Posición horizontal del cuadro de texto
top = Inches(1.5) # Posición vertical del cuadro de texto
width = Inches(9) # Ancho del cuadro de texto
height = Inches(5) # Alto del cuadro de texto
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Añadir las preguntas como párrafos al cuadro de texto
tf.add_paragraph("¿Qué necesitamos saber acerca de la audiencia?")
tf.add_paragraph("¿Por qué acuden a escucharnos o leen nuestros escritos?")
tf.add_paragraph("¿Qué esperan?")
tf.add_paragraph("¿Cuáles son sus deseos necesidades / características socioculturales?")

# Tercera diapositiva: introducción con nudo
slide = prs.slides.add_slide(slide_layout)

# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "Introducción con nudo"

# Añadir un cuadro de texto con el texto de la introducción con nudo
left = Inches(0.5) # Posición horizontal del cuadro de texto
top = Inches(1.5) # Posición vertical del cuadro de texto
width = Inches(9) # Ancho del cuadro de texto
height = Inches(5) # Alto del cuadro de texto
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Añadir el texto como un párrafo al cuadro de texto
p = tf.add_paragraph()
p.text = "Buenos días/tardes/noches. Mi nombre es _______ y estoy aquí para hablarles sobre las técnicas de expresión oral. ¿Alguna vez han tenido que hablar en público y se han sentido nerviosos, inseguros o aburridos? ¿O han asistido a alguna exposición que les ha parecido confusa, larga o irrelevante? Si es así, no se preocupen, porque hoy les voy a enseñar seis indicios que toda audiencia desea captar cuando escucha a alguien hablar. Estos indicios son señales que le damos al público para que se sienta interesado, atento y satisfecho con nuestra exposición. Si los aplicamos correctamente, podremos expresarnos oralmente de una forma clara, convincente y persuasiva. Pero antes de entrar en detalle, me gustaría hacerles una pregunta (nudo): ¿Qué es lo que más valoran ustedes cuando escuchan a alguien hablar? Piensen un momento y luego compartan sus respuestas conmigo."

# Cuarta diapositiva: desarrollo
slide = prs.slides.add_slide(slide_layout)

# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "Desarrollo"

# Añadir un cuadro de texto con el texto del desarrollo
left = Inches(0.5) # Posición horizontal del cuadro de texto
top = Inches(1.5) # Posición vertical del cuadro de texto
width = Inches(9) # Ancho del cuadro de texto
height = Inches(5) # Alto del cuadro de texto
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Añadir el texto como un párrafo al cuadro de texto
p = tf.add_paragraph()
p.text = "UNO: NO VOY A HACEROS PERDER EL TIEMPO.\nEl público se sentirá realmente molesto si tiene la sensación de que le están haciendo perder su tiempo. Es necesario dar muy pronto este indicio, a ser posible en los primeros diez segundos:\n\"Me gustaría empezar (indicio) esta breve explicación (se refuerza el indicio) preguntando a todas las personas presentes cómo creen que serían las condiciones de los cines si se bajaran los precios.\"\nDOS: SÉ QUIENES SOIS.\nEs fundamental conocer bien a la audiencia y también hacérselo saber:\n«La reducción de los precios de las entradas de cine laboral se fundamenta en los beneficios que se obtienen (indicio), que como la mayoría de los presentes (indicio reforzado), vais al cine comprobando como suben los precios día a día».\nTRES: ESTOY BIEN ORGANIZADO.\nDebemos organizar la información y, a ser posible, cómo lo estamos:\n\"En toda negociación hay dos aspectos (indicio), los intereses de los dueños de los cines y los de los cinéfilos y me gustaría hablar sobre ambos, unos minutos, antes de comentar las posibles soluciones (indicio)\".\nCUATRO: CONOZCO A FONDO EL TEMA QUE VOY A EXPONER.\nSi hemos sido presentados antes de nuestra intervención, ya se habrán destacado nuestros conocimientos y aptitudes. Pero tanto si ha sido así, como si no ha habido presentación, debemos ser nosotros quienes demostremos nuestro dominio del tema:\n\"Nos estamos reuniendo con el responsable del grupo de empresas (indicio) y os puedo asegurar que es muy receptivo a las propuestas. Así, en la próxima entrevista le daremos la documentación que hemos elaborado sobre el tema (indicio reforzado)\".\nCINCO: ESTA ES MI IDEA MÁS IMPORTANTE.\nHay que avisar cuando vamos a decir lo fundamental: \"Aunque sea lo único que nos quede claro de la charla de hoy, confío que recordaréis siempre lo que ahora os voy a comentar (indicio). En realidad se trata de la idea clave (indicio reforzado) de todo lo que he venido a exponer hoy aquí\".\nSEIS: HE TERMINADO.\n\"Antes de despedirme, y agradeciendo vuestra presencia y colaboración, me gustaría deciros...\""

# Quinta diapositiva: conclusión
slide = prs.slides.add_slide(slide_layout)

# Añadir un título a la diapositiva
title = slide.shapes.title
title.text = "Conclusión"

# Añadir un cuadro de texto con el texto de la conclusión
left = Inches(0.5) # Posición horizontal del cuadro de texto
top = Inches(1.5) # Posición vertical del cuadro de texto
width = Inches(9) # Ancho del cuadro de texto
height = Inches(5) # Alto del cuadro de texto
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame

# Añadir el texto como un párrafo al cuadro de texto
p = tf.add_paragraph()
p.text = "Hemos llegado al final de esta exposición sobre las técnicas de expresión oral. Espero que hayan disfrutado y aprendido con ella. Hemos visto los seis indicios que toda audiencia desea captar cuando escucha a alguien hablar: no hacerles perder el tiempo, saber quiénes son, estar bien organizados, conocer a fondo el tema"
ruta_guardado = 'D:/mi_presentacion.pptx'
prs.save(ruta_guardado)